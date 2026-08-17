import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

SLIDES_DIR = os.path.dirname(os.path.abspath(__file__))

# Color Palette: Modern SOC Dark / Professional Navy
NAVY = RGBColor(15, 23, 42)      # Slate 900
BLUE = RGBColor(30, 58, 138)     # Blue 900
CYAN = RGBColor(14, 165, 233)    # Sky 500
DARK_GRAY = RGBColor(51, 65, 85) # Slate 700
LIGHT_BG = RGBColor(248, 250, 252) # Slate 50
WHITE = RGBColor(255, 255, 255)
ACCENT_GREEN = RGBColor(16, 185, 129) # Emerald 500

def create_deck(title_text, subtitle_text, slides_data, output_filename):
    prs = Presentation()
    prs.slide_width = Inches(13.333) # 16:9 Widescreen
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    # --- Title Slide ---
    slide = prs.slides.add_slide(blank_layout)
    
    # Background shape
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # Title box
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"

    p2 = tf.add_paragraph()
    p2.text = subtitle_text
    p2.font.size = Pt(22)
    p2.font.color.rgb = CYAN
    p2.font.name = "Arial"
    p2.space_before = Pt(20)

    p3 = tf.add_paragraph()
    p3.text = "SD0601: SOC Analyst Workshop | 60 Hours Hands-on Curriculum"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(148, 163, 184)
    p3.space_before = Pt(30)

    # --- Content Slides ---
    for item in slides_data:
        slide = prs.slides.add_slide(blank_layout)

        # Header background banner
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.1))
        header.fill.solid()
        header.fill.fore_color.rgb = NAVY
        header.line.fill.background()

        # Header Title
        tb_head = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(11.5), Inches(0.8))
        tf_head = tb_head.text_frame
        p_head = tf_head.paragraphs[0]
        p_head.text = item["title"]
        p_head.font.size = Pt(24)
        p_head.font.bold = True
        p_head.font.color.rgb = WHITE

        # Header Subtitle / Module tag
        p_sub = tf_head.add_paragraph()
        p_sub.text = item.get("module_tag", "SD0601 SOC Analyst Curriculum")
        p_sub.font.size = Pt(12)
        p_sub.font.color.rgb = CYAN

        # Main Body - Left Column (Core Concepts)
        tb_left = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(5.7), Inches(5.5))
        tf_left = tb_left.text_frame
        tf_left.word_wrap = True

        p_l_head = tf_left.paragraphs[0]
        p_l_head.text = "📚 Core Fundamentals & Concepts"
        p_l_head.font.size = Pt(18)
        p_l_head.font.bold = True
        p_l_head.font.color.rgb = BLUE

        for pt in item.get("concepts", []):
            p = tf_left.add_paragraph()
            p.text = "• " + pt
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_GRAY
            p.space_before = Pt(8)

        # Main Body - Right Column (How to Explain & Lab Connection)
        tb_right = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(5.7), Inches(5.5))
        tf_right = tb_right.text_frame
        tf_right.word_wrap = True

        p_r_head = tf_right.paragraphs[0]
        p_r_head.text = "💡 How to Explain to Students & Lab Link"
        p_r_head.font.size = Pt(18)
        p_r_head.font.bold = True
        p_r_head.font.color.rgb = ACCENT_GREEN

        for pt in item.get("teaching_notes", []):
            p = tf_right.add_paragraph()
            p.text = "► " + pt
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_GRAY
            p.space_before = Pt(8)

        # Bottom Lab Box / Diagram Callout if present
        if "lab_callout" in item:
            callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.1), Inches(11.7), Inches(1.0))
            callout.fill.solid()
            callout.fill.fore_color.rgb = RGBColor(241, 245, 249)
            callout.line.color.rgb = CYAN

            tf_c = callout.text_frame
            tf_c.word_wrap = True
            p_c = tf_c.paragraphs[0]
            p_c.text = "🔬 Hands-On Lab Integration: " + item["lab_callout"]
            p_c.font.size = Pt(12)
            p_c.font.bold = True
            p_c.font.color.rgb = NAVY

            if "diagram_text" in item:
                p_c2 = tf_c.add_paragraph()
                p_c2.text = "📊 Visual Diagram: " + item["diagram_text"]
                p_c2.font.size = Pt(11)
                p_c2.font.color.rgb = DARK_GRAY

    output_path = os.path.join(SLIDES_DIR, output_filename)
    prs.save(output_path)
    print(f"[OK] Saved presentation: {output_path}")

# ==============================================================================
# DAY 1 SLIDES — MODULE 1: SOC FUNDAMENTALS
# ==============================================================================
day1_slides = [
    {
        "title": "Module 1: SOC Fundamentals & Roles",
        "module_tag": "Day 1 — Module 1 | TOC: Roles, Responsibilities & Blue Team Concepts",
        "concepts": [
            "Definition of Security Operations Center (SOC): Centralized unit monitoring organization's security posture.",
            "Tier 1 SOC Analyst: Initial triage, alert monitoring, true/false positive classification.",
            "Tier 2 Incident Responder: Deep investigation, correlation, containment & eradication.",
            "Tier 3 Threat Hunter / Forensic Specialist: Proactive hunting, root cause analysis, memory triage.",
            "SOC Lead / Manager: SLA tracking, escalation compliance, reporting, and operational readiness."
        ],
        "teaching_notes": [
            "Use the Hospital Emergency Room analogy: Triage Nurse (Tier 1), Doctor (Tier 2), Specialist (Tier 3).",
            "Emphasize that Tier 1 must make quick, accurate decisions within 5–10 minute SLAs.",
            "Explain key SOC metrics: MTTR (Mean Time to Respond) and MTTD (Mean Time to Detect)."
        ],
        "lab_callout": "Lab 1.1 — Mock Ticket Triage: Classify 8 SOC tickets by priority, severity, and escalation level.",
        "diagram_text": "SOC Workflow: Event -> Tier 1 Triage -> Tier 2 Investigation -> Tier 3 Root Cause / Closure"
    },
    {
        "title": "Module 1: Cyber Threat Landscape & MITRE ATT&CK",
        "module_tag": "Day 1 — Module 1 | TOC: Cyber Threat Landscape & Security Policies",
        "concepts": [
            "Cyber Threat Landscape: Phishing, Ransomware, Insider Threats, Supply Chain attacks.",
            "MITRE ATT&CK Framework: Standardized taxonomy of adversary tactics and techniques.",
            "Tactics (The 'Why'): 14 enterprise tactics representing adversary goals (e.g., Initial Access, Persistence).",
            "Techniques (The 'How'): Specific adversary actions (e.g., T1566 Spearphishing Attachment).",
            "Security Policies & Incident Response SLAs: Documented procedures governing escalation paths."
        ],
        "teaching_notes": [
            "Show students how ATT&CK replaces guesswork with a common language for attackers and defenders.",
            "Walk through an attack chain: Phishing (T1566) -> PowerShell (T1059) -> Sched Task (T1053) -> Admin User (T1136).",
            "Demonstrate ATT&CK Navigator to visually map defensive coverage."
        ],
        "lab_callout": "Lab 1.1 — ATT&CK Mapping: Map 7 simulated attacks to ATT&CK Technique IDs and build a Navigator JSON layer.",
        "diagram_text": "ATT&CK Chain Flowchart: Initial Access -> Execution -> Persistence -> Credential Access -> Lateral Movement"
    }
]

# ==============================================================================
# DAY 2 SLIDES — MODULE 2: SECURITY MONITORING & LOG MANAGEMENT
# ==============================================================================
day2_slides = [
    {
        "title": "Module 2: Log Sources & Collection Pipelines",
        "module_tag": "Day 2 — Module 2 | TOC: Log Collection, Windows, Linux & Application Logs",
        "concepts": [
            "Log Sources: Windows Event Logs, Linux syslog/auth.log, Network (Zeek/Suricata), EDR telemetry.",
            "Windows Security Events: Event ID 4624 (Logon Success), 4625 (Logon Failure), 4720 (User Created), 4698 (Task Created).",
            "Sysmon (System Monitor): Extended Windows telemetry for process creation (ID 1), network connections (ID 3), file creation (ID 11).",
            "Linux Monitoring: /var/log/auth.log (SSH logins, sudo invocations) and auditd kernel watch rules.",
            "Log Forwarders: Winlogbeat (Windows) and Filebeat (Linux) shipping logs securely to SIEM."
        ],
        "teaching_notes": [
            "Explain that default Windows logs miss process command lines — Sysmon fills this critical visibility gap.",
            "Show students how Winlogbeat & Filebeat parse raw logs into structured JSON before sending to Elasticsearch.",
            "Demonstrate Sysmon Event ID 1 parent-child process relationship (e.g., WinWord.exe spawning PowerShell.exe)."
        ],
        "lab_callout": "Lab 2.1 — Log Pipeline: Endpoint to SIEM. Configure Sysmon, Winlogbeat, Filebeat, and auditd shipping to ELK.",
        "diagram_text": "Log Pipeline: Endpoint (Sysmon/auditd) -> Log Agent (Winlogbeat/Filebeat) -> Ingestion -> SIEM Storage"
    },
    {
        "title": "Module 2: Event Correlation & Log Parsing",
        "module_tag": "Day 2 — Module 2 | TOC: Event Correlation, Log Parsing & Alert Prioritization",
        "concepts": [
            "Event Correlation: Combining multiple single log events to identify complex attack patterns.",
            "Threshold-based Correlation: e.g., 5 or more Event ID 4625 failed logins from same IP within 5 minutes = Brute Force.",
            "Sequence-based Correlation: Event ID 4720 (User Created) followed immediately by 4732 (Added to Admin Group).",
            "Log Parsing: Extracting key-value pairs (src_ip, user, process_name) from unstructured log lines.",
            "Alert Generation & Prioritization: Severity scoring (Low, Medium, High, Critical) based on asset criticality."
        ],
        "teaching_notes": [
            "Emphasize that 1 failed login is normal noise; 50 failed logins in 60 seconds is a high-confidence alert.",
            "Teach students how correlation rules reduce alert fatigue for Tier 1 analysts.",
            "Walk through writing a basic correlation threshold logic rule."
        ],
        "lab_callout": "Lab 2.1 — Correlation Rule Build: Deploy 5+ SSH failed login correlation rule in Kibana/Elasticsearch.",
        "diagram_text": "Correlation Engine: Raw Events -> Parsing/Normalizing -> Rule Matching -> Threshold Evaluation -> Alert"
    }
]

# ==============================================================================
# DAY 3 SLIDES — MODULE 3: SIEM & LOG ANALYSIS (SPLUNK & ELK)
# ==============================================================================
day3_slides = [
    {
        "title": "Module 3: SIEM Architecture & ELK Fundamentals",
        "module_tag": "Day 3 — Module 3 | TOC: SIEM Architecture, ELK Stack & Kibana Searching",
        "concepts": [
            "SIEM Core Functions: Ingestion, Indexing, Storage, Search, Correlation, Alerting, Dashboarding.",
            "ELK Stack Components: Elasticsearch (search engine), Logstash (log transformer), Kibana (visualization UI).",
            "Kibana KQL (Kibana Query Language): Filtering syntax for fast event searching (e.g., event.code:4625 AND event.outcome:failure).",
            "Data Views / Index Patterns: Mapping index name patterns (winlogbeat-*, filebeat-*) to Kibana Discover.",
            "Kibana Dashboards: Building visual grids with bar charts, pie charts, data tables, and metric indicators."
        ],
        "teaching_notes": [
            "Compare SIEM to a Google Search engine specifically for enterprise security logs.",
            "Demonstrate KQL filtering in Kibana Discover live: searching for specific users, Event IDs, and source IPs.",
            "Show how to build a 3-panel dashboard: Top Failed Users, Failed Login Trend Chart, Source IP Table."
        ],
        "lab_callout": "Lab 3.1 — Kibana Build-Out: Write 5 KQL queries, configure Kibana data views, and build a Failed Login Dashboard.",
        "diagram_text": "ELK SIEM Architecture: Agents -> Beats -> Elasticsearch (Index) -> Kibana (Search & Dashboards)"
    },
    {
        "title": "Module 3: Splunk Fundamentals & SPL Searching",
        "module_tag": "Day 3 — Module 3 | TOC: Introduction to Splunk, SPL Search Queries & Alerts",
        "concepts": [
            "Splunk Architecture: Forwarder -> Indexer -> Search Head.",
            "Splunk SPL (Search Processing Language): Pipeline-based search syntax using standard pipes (|).",
            "SPL Query 1 (Brute Force): index=main sourcetype='WinEventLog:Security' EventCode=4625 | stats count by src_ip | where count > 5.",
            "SPL Query 2 (Port Scan): index=zeek sourcetype='zeek:conn' | stats dc(dest_port) as unique_ports by src_ip | where unique_ports > 20.",
            "SPL Query 3 (C2 Beaconing): Calculating regularity using avg() and stdev() over duration.",
            "Splunk Alerting: Saving SPL queries as scheduled alerts triggered on search threshold conditions."
        ],
        "teaching_notes": [
            "Highlight the difference between Kibana KQL (filter-first) and Splunk SPL (pipeline statistics engine).",
            "Teach students how the pipe character '|' passes search results to transform functions like stats, eval, where, table.",
            "Walk through the beaconing detection formula: low standard deviation in connection interval = automated C2 beacon."
        ],
        "lab_callout": "Lab 3.1 — Splunk Build-Out: Write SPL for brute force, port scans, and C2 beaconing in Splunk 9.2 Free.",
        "diagram_text": "Splunk Search Pipeline: index search | stats aggregation | eval transformation | where filter | table output"
    }
]

# ==============================================================================
# DAY 4 SLIDES — MODULE 4: THREAT DETECTION & INCIDENT RESPONSE
# ==============================================================================
day4_slides = [
    {
        "title": "Module 4: Indicators of Compromise & YARA Analysis",
        "module_tag": "Day 4 — Module 4 | TOC: IOCs, Malware Detection Basics & YARA Rules",
        "concepts": [
            "Indicators of Compromise (IOCs): Artifacts observed on network/host indicating breach (Hashes, IPs, Domains, URLs).",
            "Pyramid of Pain: Hashes (Easy to change) -> IPs -> Domains -> Host Artifacts -> TTPs (Hardest for attacker to change).",
            "Malware Detection Basics: Static analysis (strings, hashes, headers) vs Dynamic analysis (sandboxing, behavioral execution).",
            "YARA Rules: Pattern matching Swiss army knife for malware identification based on text, hex strings, and logical conditions.",
            "YARA Rule Structure: meta (author, description, ATT&CK tag), strings ($string1, $hex1), condition ($string1 and $hex1)."
        ],
        "teaching_notes": [
            "Explain Bianca Rivera's Pyramid of Pain — why relying only on file hashes is insufficient for modern SOC defense.",
            "Show students how YARA rules allow SOC analysts to detect malware variants even when the hash changes.",
            "Walk through writing a YARA rule for detecting macro droppers and encoded PowerShell strings."
        ],
        "lab_callout": "Lab 4.1 — YARA Rule Engineering: Write and test suspicious_rule.yar against lab sample files.",
        "diagram_text": "Pyramid of Pain: Hashes (Base) -> Domain Names -> Network Artifacts -> Tools -> TTPs (Apex)"
    },
    {
        "title": "Module 4: Incident Response Lifecycle & Case Management",
        "module_tag": "Day 4 — Module 4 | TOC: NIST IR Lifecycle, Threat Intel (MISP) & Case Management (TheHive)",
        "concepts": [
            "NIST SP 800-61 IR Lifecycle: 1. Preparation -> 2. Detection & Analysis -> 3. Containment, Eradication & Recovery -> 4. Post-Incident Activity.",
            "Threat Intelligence Platforms (MISP): Sharing, storing, and enriching structured IOCs with TLP markings (TLP:AMBER) and ATT&CK galaxy tags.",
            "Case Management Platforms (TheHive 5): Organising security incidents, assigning observables, tasks, and tracking analysis timelines.",
            "Cortex Analyzer Engine: Automated enrichment of IOC observables via external APIs (VirusTotal, AbuseIPDB).",
            "TLP Markings: Traffic Light Protocol (TLP:RED, TLP:AMBER, TLP:GREEN, TLP:CLEAR) controlling threat data sharing."
        ],
        "teaching_notes": [
            "Walk through a full case lifecycle: Alert fires -> Created in TheHive -> Observables enriched via Cortex -> Exported to MISP.",
            "Emphasize that Containment must happen BEFORE Eradication (e.g., isolate host before killing malicious process).",
            "Explain how Threat Intelligence sharing prevents other organizations from falling victim to the same threat campaign."
        ],
        "lab_callout": "Lab 4.1 — Full Case Lifecycle: Create MISP event, enrich IOCs, run Cortex analyzers, and progress a case in TheHive 5.",
        "diagram_text": "TheHive IR Workflow: Alert Ingestion -> Observable Enrichment (Cortex) -> Analysis Tasks -> Containment -> MISP Export"
    }
]

# ==============================================================================
# DAY 5 SLIDES — MODULE 5 & 6: ENDPOINT HUNTING, FORENSICS & LIVE DRILL
# ==============================================================================
day5_slides = [
    {
        "title": "Module 5: EDR, Threat Hunting & Memory Forensics",
        "module_tag": "Day 5 — Module 5 | TOC: EDR (Wazuh), Threat Hunting & Volatility 3 Memory Triage",
        "concepts": [
            "Endpoint Detection & Response (EDR): Continuous monitoring, telemetry collection, and automated active response on endpoints.",
            "Wazuh Open Source EDR: Manager, Indexer, Dashboard architecture with rule-based decoders (e.g., detecting Event 4720/4698).",
            "Threat Hunting Methodology: Proactive hypothesis-driven searching for undetected threats (Hypothesis -> Query -> Finding -> Documentation).",
            "Memory Forensics (Volatility 3): Analyzing RAM dumps to detect malware, injected code, hidden processes, and unlinked sockets.",
            "Volatility 3 Commands: windows.info, windows.pslist (process list), windows.pstree (process parent-child tree), windows.netscan (sockets)."
        ],
        "teaching_notes": [
            "Differentiate SIEM (centralized log aggregation) vs EDR (deep endpoint activity & active response agent).",
            "Teach the threat hunting mindset: 'Assuming a breach has occurred, how would I find an Office process spawning PowerShell?'",
            "Demonstrate Volatility 3: finding an anomalous svchost.exe whose parent is NOT services.exe."
        ],
        "lab_callout": "Lab 5.1 — EDR & Memory Triage: Enroll Wazuh agents, execute threat hunt hypothesis, run Volatility 3 commands on victim.mem.",
        "diagram_text": "Threat Hunting Cycle: Formulate Hypothesis -> Develop Query -> Analyze Telemetry -> Uncover Finding -> Document & Mitigate"
    },
    {
        "title": "Module 6: Live-Fire Drill, Capstone & Interview Prep",
        "module_tag": "Day 5 — Module 6 | TOC: Atomic Red Team Live Drill, Capstone Project & Analyst Interview Prep",
        "concepts": [
            "Atomic Red Team: Open-source framework for executing small, controlled, mapping-aligned attack tests (T1110, T1053, T1059).",
            "Live-Fire Drill SLA Tracking: Tier 1 Detection SLA (<= 5 min), Triage SLA (<= 2 min), Escalation SLA (<= 5 min).",
            "Operation SilentLedger Capstone: Multi-stage breach investigation across Phishing, Execution, Persistence, Lateral Movement, and Exfiltration.",
            "SOC Analyst Documentation: Professional incident timeline recording attack timestamps, detection timestamps, and containment steps.",
            "SOC Analyst Technical Interview Preparation: Walkthrough of technical questions (Logon Event IDs, TCP 3-way handshake, SIEM queries, Pyramid of Pain)."
        ],
        "teaching_notes": [
            "Pair students into Attacker / Defender roles for the 45-minute live-fire drill.",
            "Stress that SLAs simulate real-world SOC pressure — speed and accuracy are equally important.",
            "Review the top 10 SOC Analyst technical interview questions and guide students on structuring clear, confident answers."
        ],
        "lab_callout": "Lab 6.1 & Capstone: Execute Atomic Red Team tests, triage alerts in real-time, complete Operation SilentLedger, prepare final report.",
        "diagram_text": "Live-Fire Drill Flow: Attacker Executes Test -> Telemetry Generated -> SIEM Alert Fires -> Defender Triages -> Case Created in SLA"
    }
]

# Generate all 5 individual day presentations plus a master presentation
create_deck("SD0601 — Day 1: SOC Fundamentals & ATT&CK", "Module 1: Roles, Responsibilities, Threat Landscape & ATT&CK Taxonomy", day1_slides, "SD0601_Day1_SOC_Fundamentals.pptx")
create_deck("SD0601 — Day 2: Security Monitoring & Log Management", "Module 2: Log Sources, Sysmon, Filebeat, Auditd & Event Correlation", day2_slides, "SD0601_Day2_Log_Management.pptx")
create_deck("SD0601 — Day 3: SIEM & Log Analysis (Splunk & ELK)", "Module 3: SIEM Architecture, KQL, SPL Search Queries & Dashboards", day3_slides, "SD0601_Day3_SIEM_Splunk_ELK.pptx")
create_deck("SD0601 — Day 4: Threat Detection & Incident Response", "Module 4: IOCs, YARA Rule Engineering, MISP, Cortex & TheHive 5", day4_slides, "SD0601_Day4_Threat_Detection_IR.pptx")
create_deck("SD0601 — Day 5: Endpoint Hunting, Forensics & Live Drill", "Module 5 & 6: Wazuh EDR, Volatility 3, Atomic Red Team & Capstone", day5_slides, "SD0601_Day5_Endpoint_Forensics_LiveDrill.pptx")

# Master Deck
master_slides = day1_slides + day2_slides + day3_slides + day4_slides + day5_slides
create_deck("SD0601 — Complete Master Presentation Deck", "SOC Analyst: Security Monitoring and Incident Response (60 Hours)", master_slides, "SD0601_Master_Presentation_Deck.pptx")

print("All PowerPoint (.pptx) decks generated successfully in slides/")
